# -*- coding: utf-8 -*-
"""
Created on Thu Oct 28 17:50:08 2021

@author: SNASONOV
"""

import os, subprocess, datetime
import geopandas as gpd
import pandas as pd
from osgeo import gdal

from pptx import Presentation
from pptx.util import Cm, Inches
from pptx.util import Pt


def getfiles(d,ext):
    paths = []
    for file in os.listdir(d):
        if file.endswith(ext):
            paths.append(os.path.join(d, file))
    return(paths) 

mydir = r'C:\Data\Burn_Severity\ElephantHill\qa'
pre_folder = os.path.join(mydir,'pre')
pre = getfiles(pre_folder,'.png')

post_folder = os.path.join(mydir,'post')
post = getfiles(post_folder,'.png')

barc_folder = os.path.join(mydir,'barc')
barc = getfiles(barc_folder,'.png')

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

title_only_slide_layout = prs.slide_layouts[5]

for i,j,k in zip(pre,post,barc):
    slide = prs.slides.add_slide(title_only_slide_layout)
    title = slide.shapes.title
    
    name = i.rsplit('\\')[-1]
    name2 = name.rsplit('_')[0]
    
    title.text = name2
    #title.font.size = Pt(20)
    
    slide.shapes.add_picture(
        i, left=Cm(0.44), top=Cm(3.7), width=Cm(12.24), height=None
    )
    slide.shapes.add_picture(
        j, left=Cm(12.78), top=Cm(3.7), width=Cm(12.24), height=None
    )
    slide.shapes.add_picture(
        k, left=Cm(25.02), top=Cm(3.7), width=Cm(12.24), height=None
    )
    
prs.save(os.path.join(mydir,'qa.pptx'))