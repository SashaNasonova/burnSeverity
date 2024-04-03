# -*- coding: utf-8 -*-
"""
Created on Thu Apr 27 16:44:26 2023

@author: snasonov
"""

import os
import pandas as pd
import geopandas as gpd
from osgeo import gdal
import rasterio
from rasterio.plot import show
import matplotlib
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
import numpy as np
from pathlib import Path

import cartopy.crs as ccrs
import cartopy.feature as cfeature

from pptx import Presentation
from pptx.util import Cm, Inches
from pptx.util import Pt


def getfiles(d,ext):
    paths = []
    for file in os.listdir(d):
        if file.endswith(ext):
            paths.append(os.path.join(d, file))
    return(paths) 


def ql_3band(outshp,imgpath,outpath): 
    # Get fire vector
    dfs_sub = gpd.read_file(outshp)
    
    # Read 8-bit image
    src = rasterio.open(imgpath)
    
    # Create pre-fire quicklook image with burned area boundary overlay
    fig, ax = plt.subplots(figsize=(8, 8)) #10,13
    base = show(src,ax=ax)
    dfs_sub.plot(ax=base, edgecolor='purple', facecolor='none',linewidth=2)
    ax.axis('off')
    plt.savefig(outpath, bbox_inches='tight', dpi=300)
    
    plt.close()
    src = None

def generate_legend_labels(arr):
    #sort array
    arr.sort()
    
    default_labels = {
        0: {"black":"Unknown"},
        1: {"gray":"Unburned"},
        2: {"yellow":"Low"},
        3: {"orange":"Medium"},
        4: {"red":"High"}
    }
    
    legend_labels = {}
    
    for value in arr:
        #print(default_labels[value]) #for debug
        legend_labels.update(default_labels[value])
    
    return(legend_labels)


def ql_barc(outshp,barcpath,imgpath,outpath):
    dfs_sub = gpd.read_file(outshp)
    
    # Read in BARC classification
    src = rasterio.open(barcpath)
    unique_values = np.unique(src.read(1))
    
    # Read in post-fire image 
    src1 = rasterio.open(imgpath)
    
    #drop no data 9 
    arr = unique_values[unique_values!=9]
    
    # Create cmap dictionary
    colormap_dict = {0: 'black', 
                     1: 'gray', 
                     2: 'yellow',
                     3: 'orange', 
                     4: 'red'}
    
    colormap_unique = [colormap_dict[value] for value in arr]
    colormap = matplotlib.colors.ListedColormap(colormap_unique)
    
    # Create pre-fire quicklook image with burned area boundary overlay
    fig, ax = plt.subplots(figsize=(8, 8))
    background = show((src1,1),ax=ax,cmap='gray')
    base = show(src,cmap=colormap,ax=ax)
    dfs_sub.plot(ax=base, edgecolor='purple', facecolor='none',linewidth=2)
    ax.axis('off')
    
    #add legend
    # Add a legend for labels
    # legend_labels = {"black": "Unknown", 
    #                  "gray": "Unburned",
    #                  "yellow":"Low",
    #                  "orange":"Medium",
    #                  "red": "High"}
    
    legend_labels = generate_legend_labels(arr)
    
    patches = [Patch(color=color, label=label)
               for color, label in legend_labels.items()]
    
    ax.legend(handles=patches,
              bbox_to_anchor=(1.05, 1),loc='upper left',
              borderaxespad=0., facecolor="white")
    
    
    plt.savefig(outpath, bbox_inches='tight', dpi=300)
    
    #plt.savefig(outpath_pdf, format="pdf", bbox_inches="tight")
    
    plt.close()
    src = None
    src1 = None
    
def create_ppt(pptpath):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    prs.save(pptpath)

def add_slide(pptpath,i,j,k,l,df):
    prs = Presentation(pptpath)
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    title = slide.shapes.title
    title.width = Cm(35)
    title.height = Cm(3.3)
    title.top = Cm(0)
    title.left = Cm(0)
    
    #Add title
    name = Path(k).stem
    title.text = name
    
    #First image (pre-fire)
    slide.shapes.add_picture(
        i, left=Cm(0.44), top=Cm(8.2), width=Cm(12.24), height=None
    )
    #Second image (post_fire)
    slide.shapes.add_picture(
        j, left=Cm(12.78), top=Cm(8.2), width=Cm(12.24), height=None
    )
    
    #get image height
    picture = slide.shapes[1] #second shape pre-img, first is the title
    height_cm = picture.height.cm
    #print(height_cm)
    
    #Third image (barc)
    slide.shapes.add_picture(
        k, left=Cm(25.02), top=Cm(8.2), width=None, height=Cm(height_cm)
    )
    
    # #Third image (barc)
    # slide.shapes.add_picture(
    #     k, left=Cm(25.02), top=Cm(8.2), width=Cm(12.24), height=None
    # )
    
    #Fourth image (location map)
    slide.shapes.add_picture(
        l, left=Cm(32.5), top=Cm(0), width=None, height=Cm(7)
    )
    
    #add table
    rows, cols = df.shape
    left = Cm(0.9)
    top = Cm(3.2)
    width = Cm(16.3)
    height = Cm(3.75)

    #add table
    shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = shape.table

    #assign table style    
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}'
    tbl[0][-1].text = style_id
 
    
    # Set column names
    for col, column_name in enumerate(df.columns):
        cell = table.cell(0, col)
        cell.text = column_name
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # Populate data
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row + 1, col)
            cell.text = str(df.iloc[row, col])
    
    # Change font size
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
    
    prs.save(pptpath) #this overwrites
    print('Presentation saved')
    
def ql_3band_batch(folder,outshp,outfolder):    
    imglist = getfiles(folder,'.tif')
    for imgpath in imglist:
        name = Path(imgpath).stem + '.png'
        outpath = os.path.join(outfolder,name)
        ql_3band(outshp,imgpath,outpath)

def add_slides_batch(img_folder,pptpath):
    imglist = getfiles(img_folder,'.png')
    
    prs = Presentation(pptpath)
    title_only_slide_layout = prs.slide_layouts[5]
    
    for i in imglist:
        print(i)
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        title.width = Cm(35)
        title.height = Cm(3.3)
        title.top = Cm(0)
        title.left = Cm(0)
        
        #Add title
        name = Path(i).stem
        title.text = name
        
        slide.shapes.add_picture(
            i, left=Cm(12.78), top=Cm(3.7), width=Cm(12.24), height=None
        )
    prs.save(pptpath)
    
def inset_map(bc_path,fire_perim,outpath):
    # Define the bounding box for British Columbia (lonmin, lonmax, latmin, latmax)
    bbox = [-139, -114.75, 47.5, 60]
    
    
    # Coordinates of cities/towns (latitude, longitude)
    cities = {
        'Vancouver': (49.2827, -123.1207),
        'Kamloops': (50.6761, -120.3408),
        'Prince George':(53.9170,-122.7494),
        'Fort St John': (56.2464,-120.8476),
        'Prince Rupert': (54.3125,-130.3054),
        'Williams Lake': (52.1284,-122.1302)
        
    }
    
    # Create a map 
    plt.figure(figsize=(4,4))
    ax = plt.axes(projection=ccrs.AlbersEqualArea(central_longitude=-126, central_latitude=54))
    ax.set_extent(bbox, crs=ccrs.PlateCarree())
    
    # Add land and coastline features
    ax.add_feature(cfeature.LAND, edgecolor='black', facecolor='lightgrey',linewidth=0.1)
    ax.add_feature(cfeature.COASTLINE, linewidth=0.1)
    
    # Load the British Columbia boundary data using GeoPandas
    ## Thank you GeoBC! https://catalogue.data.gov.bc.ca/dataset/province-of-british-columbia-boundary-terrestrial
    bc_boundary = gpd.read_file(bc_path)
    
    # Plot the British Columbia boundary using Cartopy's geopandas tools
    ax.add_geometries(bc_boundary['geometry'], crs=ccrs.PlateCarree(), edgecolor='black', facecolor='green',linewidth=0.1)
    
    # Add gridlines
    ax.gridlines(draw_labels=True, linestyle='--', color='grey')
    
    # Add fire boundary
    #fire_perim = r"C:\Data\Burn_Severity\same_year_2023\tumbler_ridge\vectors\G70645_temp_gcs.shp"
    fire_boundary = gpd.read_file(fire_perim)
    fire_boundary["centroid"] = fire_boundary["geometry"].centroid
    lat = fire_boundary["centroid"].y
    lon = fire_boundary["centroid"].x
    ax.plot(lon,lat, marker='*', color='purple', markersize=5, transform=ccrs.PlateCarree())
    ax.text(lon + 0.5, lat, 'Fire Location', color='purple', fontsize=8, transform=ccrs.PlateCarree())
    
    # Add cities/towns
    for city, (lat, lon) in cities.items():
        ax.plot(lon, lat, marker='o', color='black', markersize=3, transform=ccrs.PlateCarree())
        ax.text(lon + 0.5, lat, city, color='black', fontsize=7, transform=ccrs.PlateCarree())
    
    #outpath = r'C:\Data\Burn_Severity\code_test\map.png'
    plt.savefig(outpath, bbox_inches='tight', dpi=300)

def zonal_barc(barcpath,firepath,outpath):
    
    def burnsev_name(x):
        if x == 0:
            return('Unknown')
        elif x == 1:
            return('Unburned')
        elif x == 2:
            return('Low')
        elif x == 3: 
            return('Medium')
        elif x == 4: 
            return('High')
        else:
            print("Wrong burn severity value!")
        
    
    #need to clip to fire perimeter a little tighter than what gee outputs
    basename = barcpath[:-4] 
    barcpath_clip = basename + '_clip.tif'
    
    gdal.Warp(barcpath_clip,barcpath,cutlineDSName=firepath,cropToCutline=True)
    
    dat = gdal.Open(barcpath_clip)
    band = dat.GetRasterBand(1).ReadAsArray()
    
    x1,px,x2,x3,x4,x5 = dat.GetGeoTransform()
      
    nodatavalue = 9
    vals, counts = np.unique(band[band != nodatavalue], return_counts=True)
    
    df = pd.DataFrame(
        {'class':vals,
         'px_count':counts})
    
    
    a = df['px_count'].sum()
    df['perc' ] = (df['px_count'] / a)*100
    df['area_m2'] = df['px_count']*(px*px)
    df['area_ha'] = df['area_m2']*0.0001
    
    #round to 1 decimal place
    df = df.round(1)
    df['burn_sev'] = df['class'].apply(burnsev_name)
    
    #reorder columns
    df = df[['class','burn_sev','px_count','area_m2','area_ha','perc']]
    
    #print(df) #debug
    df.to_csv(outpath)
    return(df)        
             
